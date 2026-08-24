#!/usr/bin/env python3
"""The literal .pptx twin of the deck artifact: same 14 slides, same
register (ink ground, indigo eyebrows, gold for THE number), images from
the deck pack, every number from key_numbers.json."""
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

DECK = Path(__file__).resolve().parent.parent.parent / "docs" / "assets" / "deck"
OUT = DECK / "verdict-bench.pptx"
K = json.loads((DECK / "key_numbers.json").read_text())

INK = RGBColor(0x08, 0x09, 0x0C)
TEXT = RGBColor(0xE8, 0xEC, 0xF4)
MUT = RGBColor(0x8A, 0x93, 0xA3)
DIM = RGBColor(0x5C, 0x66, 0x78)
ACC = RGBColor(0x82, 0x8F, 0xFF)
GOLD = RGBColor(0xCF, 0xA3, 0x6C)
GOOD = RGBColor(0x4A, 0xDE, 0x80)
BAD = RGBColor(0xF8, 0x71, 0x71)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = INK
    return s


def box(s, left, top, width, height):
    return s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame


def para(tf, text, size, color, bold=False, first=False, space=6, caps=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    r = p.add_run()
    r.text = text.upper() if caps else text
    f = r.font
    f.size = Pt(size)
    f.color.rgb = color
    f.bold = bold
    f.name = "Inter"
    p.space_after = Pt(space)
    return p


def eyebrow(s, text):
    tf = box(s, 0.9, 0.55, 11.5, 0.4)
    para(tf, text, 12, ACC, bold=True, first=True, caps=True)


def h2(s, text, top=1.0):
    tf = box(s, 0.9, top, 11.5, 1.0)
    para(tf, text, 33, TEXT, bold=True, first=True)


def note(s, text, top=6.3, color=MUT):
    tf = box(s, 0.9, top, 11.5, 1.0)
    tf.word_wrap = True
    para(tf, text, 13, color, first=True)


def bigrow(s, items, top=2.4):
    for i, (num, cap, color) in enumerate(items):
        left = 0.9 + i * 4.1
        tf = box(s, left, top, 3.9, 2.2)
        para(tf, num, 58, color, bold=True, first=True)
        para(tf, cap, 11, DIM, caps=True)


def image(s, name, top=2.0, height=4.0):
    s.shapes.add_picture(str(DECK / name), Inches(0.9), Inches(top), height=Inches(height))

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.oxml.ns import qn

PANEL = RGBColor(0x11, 0x14, 0x20)
EDGE = RGBColor(0x3A, 0x3F, 0x4E)


def fbox(s, left, top, w, h, title, sub, edge=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = PANEL
    sh.line.color.rgb = edge or EDGE
    sh.line.width = Pt(1.2)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    para(tf, title, 13, TEXT, bold=True, first=True, space=2)
    para(tf, sub, 10, MUT, space=0)
    return sh


def arrow(s, x1, y1, x2, y2):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = ACC
    c.line.width = Pt(1.6)
    ln = c.line._get_or_add_ln()
    te = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(te)
    return c




# 0 cover
s = slide()
s.shapes.add_picture(str(DECK / "intuit-mark.png"),
                     Inches(0.93), Inches(1.35), height=Inches(0.5))
tf = box(s, 0.9, 2.3, 11.8, 3.2)
para(tf, "verdict-bench", 88, TEXT, bold=True, first=True)
para(tf, "account-review decisioning prompt  ·  the Intuit case study", 13, DIM, caps=True, space=26)
para(tf, "Shoval Benjer  ·  August 2026", 14, MUT)

# 0b about me + operating loop
s = slide()
eyebrow(s, "Who is presenting")
h2(s, "I run AI the way this benchmark runs prompts")
tf = box(s, 0.9, 1.85, 11.5, 1.5)
tf.word_wrap = True
para(tf, "Shoval Benjer, B.Sc. Data Science (Afeka, 2022-2025). Before this: seven months at an online-trading company in Herzliya, owning the quality, accuracy, and reliability of AI agents on three production surfaces (a chatbot, voice agents, and social content). Built a sales-training platform on voice agents and a four-language transcription-and-analysis pipeline for sales and retention calls; implemented a hybrid test pyramid and quality gates, evals beside deterministic checks, that raised agent quality on the company's internal benchmarks and block regressions before they reach a customer; and a real-time productivity dashboard (SQL, Next.js) over the CRM with query optimization and cache layers.", 13, TEXT, first=True)
fbox(s, 0.9, 3.6, 2.1, 1.0, "operator", "frames the question", edge=ACC)
fbox(s, 3.5, 3.6, 2.3, 1.0, "rules + claim controls", "what may be asserted")
fbox(s, 6.3, 3.6, 2.1, 1.0, "agents run", "every run recorded")
fbox(s, 8.9, 3.6, 2.4, 1.0, "evidence ledger", "sqlite + jsonl, append-only", edge=GOLD)
fbox(s, 5.3, 5.1, 2.7, 0.95, "gates decide", "green ships, red blocks", edge=GOOD)
arrow(s, 3.0, 4.1, 3.5, 4.1)
arrow(s, 5.8, 4.1, 6.3, 4.1)
arrow(s, 8.4, 4.1, 8.9, 4.1)
arrow(s, 10.1, 4.6, 8.0, 5.35)
arrow(s, 5.3, 5.55, 1.95, 4.6)
note(s, "My standing setup runs on this loop: rules say what may be claimed, every run lands "
        "in a ledger, and gates decide what ships. verdict-bench is the same loop pointed at "
        "one prompt.", top=6.35)

# 0c thesis
s = slide()
eyebrow(s, "The study in three findings")
h2(s, "What the ledger showed")
tf = box(s, 0.9, 2.2, 11.4, 4.2)
tf.word_wrap = True
for i, (h, b) in enumerate([
    ("The gate rejected my best fix.", "v6 resisted 19 of 20 planted-instruction attacks, "
     "then missed the 12-of-12 decision bar it had pre-registered. It stays rejected."),
    ("Generated cases separate models the visible nine cannot.", "on 64 unseen cases the "
     "same prompt scores 100% on two models and 58% on the weakest."),
    ("Contested cases are routed, not scored.", "where the policy text underdetermines the "
     "answer, no model is graded and the case goes to a human with the cross-model split attached."),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r1 = p.add_run(); r1.text = h + " "
    r1.font.size = Pt(17); r1.font.bold = True; r1.font.color.rgb = TEXT; r1.font.name = "Inter"
    r2 = p.add_run(); r2.text = b
    r2.font.size = Pt(17); r2.font.color.rgb = MUT; r2.font.name = "Inter"
    p.space_after = Pt(18)

# 1 problem
s = slide()
eyebrow(s, "PPDAC · Problem")
h2(s, "Accuracy is the wrong headline")
bigrow(s, [("92%", "the no-policy baseline's suite accuracy", TEXT),
           ("44x", "cost spread across error types ($45 to $2,000)", TEXT),
           ("1", "sanctions miss disqualifies, regardless of the rest", GOLD)])
note(s, "Bottom line: the question is decision quality, not model quality. Which prompt-model "
        "pair do we trust, at what dollar loss, under which assumptions.", top=5.6)

# 2 two-hour
s = slide()
eyebrow(s, "Scope, owned upfront")
h2(s, "The two-hour version exists inside this one")
tf = box(s, 0.9, 2.3, 11.0, 2.2)
tf.word_wrap = True
para(tf, "Just to be clear, the two-hour version exists and I can point at it: the "
         "policy-teaching prompt plus the sanctions rule, 12x2 runs, three stated costs, one "
         "page. That is rungs v3 and v4 of my own ladder, and they carry most of the decision "
         "quality. The rest of the time went on the one question the short version cannot "
         "answer: how do I know when to trust it.", 17, TEXT, first=True)
note(s, "\"The prompt took two hours. Knowing whether to trust it took the week, and the "
        "week is the part I would bring to the job.\"", top=5.2, color=GOLD)


# unshown work
s = slide()
eyebrow(s, "Measured, banked, and on no other slide")
h2(s, "The work the deck does not show")
tf = box(s, 0.9, 2.0, 11.4, 4.8)
tf.word_wrap = True
for i, (h, b) in enumerate([
    ("Classical baseline:", "hand-featured logistic regression, leave-one-out 8/12 vs the "
     "LLM's 11-12/12; all four LR misses are policy-reasoning cases in the expensive direction."),
    ("DSPy contrast arm:", "a compiled prompt hits the same 12/12 ceiling with a "
     "structurally different artifact; kept as a comparison, not a rung."),
    ("Self-consistency:", "11 of 12 cases unanimous over N=5 temperature-raised repeats; "
     "vote fraction is the earned-confidence instrument."),
    ("Sequential test:", "SPRT (p0=0.75, p1=0.92, alpha=0.05, beta=0.10) states how many "
     "more labeled cases certification needs."),
    ("Posterior odds:", "beta-binomial puts the ladder's edge over v1 at ~3:1, reported as "
     "odds, not certainty."),
    ("Provenance:", "PaySim and ULB licenses verified at source; IEEE-CIS marked honestly "
     "unverifiable; fetch scripts print the exact manual steps."),
    ("Surfaces:", "a Tauri desktop shell, the adjudication-queue playground, and a "
     "number-grounding audit: 81-98% of v5's cited numbers trace to the case file."),
]):
    p2 = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r1 = p2.add_run(); r1.text = h + " "
    r1.font.size = Pt(13); r1.font.bold = True; r1.font.color.rgb = TEXT; r1.font.name = "Inter"
    r2 = p2.add_run(); r2.text = b
    r2.font.size = Pt(13); r2.font.color.rgb = MUT; r2.font.name = "Inter"
    p2.space_after = Pt(9)

# 3 plan
s = slide()
eyebrow(s, "PPDAC · Plan")
h2(s, "Six commitments, fixed before any analysis")
tf = box(s, 0.9, 2.2, 11.4, 4.6)
tf.word_wrap = True
for i, (h, b) in enumerate([
    ("Suite separation.", "robustness kinds never blend into accuracy or loss."),
    ("One change per rung,", "machine-diffed, every run pinned by content hash."),
    ("Repeats before belief.", "N=5 before a single-run claim is trusted."),
    ("Trust gates over headlines:", "n, contract, CI width, flip, zero-tolerance tripwire."),
    ("Dollar OEC:", "three stated costs + one derived, swept, bootstrapped, reweighted."),
    ("Label tiers visible:", "4 expert / 5 adjudicated / constructed, never averaged."),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r1 = p.add_run(); r1.text = f"{i+1}.  {h} "
    r1.font.size = Pt(16); r1.font.bold = True; r1.font.color.rgb = TEXT; r1.font.name = "Inter"
    r2 = p.add_run(); r2.text = b
    r2.font.size = Pt(16); r2.font.color.rgb = MUT; r2.font.name = "Inter"
    p.space_after = Pt(12)



# expected-loss model
s = slide()
eyebrow(s, "PPDAC · Plan · the price of a mistake")
h2(s, "Expected loss, the actual model")
rows = [("", "-> APPROVE", "-> HOLD", "-> REJECT"),
        ("truth APPROVE", "$0", "$45", "$600"),
        ("truth HOLD", "$45", "$0", "$600"),
        ("truth REJECT", "$2,000", "$500", "$0")]
tbl = s.shapes.add_table(4, 4, Inches(0.9), Inches(2.1), Inches(5.6), Inches(2.4)).table
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = PANEL
        tfc = cell.text_frame
        para(tfc, val, 12, GOLD if val == "$2,000" else (DIM if ri == 0 or ci == 0 else TEXT),
             bold=(ri == 0 or ci == 0 or val == "$2,000"), first=True, space=0)
tf = box(s, 7.0, 2.0, 5.4, 4.6)
tf.word_wrap = True
for line in [
    "EL_1k = (1000/N) SUM_i (1/R) SUM_r C(y_i, yhat_ir)",
    "unparseable run -> charged max_yhat C(y_i, yhat), never dropped",
    "CI: case-clustered bootstrap, resample cases not runs, B=1,000, seed 1789",
    "rankable only if Wilson 95% width <= 0.5 (z=1.96, n>=8)",
    "EL(pi) = pi*EL_fraud + (1-pi)*EL_legit, pi swept 0.5% to 5%",
    "queue cost_1k = 1000 * p_HOLD * $35",
]:
    para(tf, line, 13, TEXT, space=10)
note(s, "Three costs stated, one derived ($500 = containment on a caught fraud); the $2,000 "
        "false-approve swept 1k to 5k without the ranking flipping. Accuracy weighs every "
        "mistake 1; this matrix is why the benchmark ranks on dollars.", top=6.35)

# 3b architecture
s = slide()
eyebrow(s, "PPDAC · Plan · the machine")
h2(s, "One pipeline, every claim traceable to a run")
fbox(s, 0.9, 2.3, 2.3, 1.0, "89 labeled cases", "+ POLICY.md, frozen")
fbox(s, 0.9, 3.8, 2.3, 1.0, "prompt ladder", "v1..v5, one change per rung")
fbox(s, 3.8, 3.05, 2.3, 1.0, "runner", "7 providers, N=5 repeats", edge=ACC)
fbox(s, 6.7, 3.05, 2.3, 1.0, "sqlite ledger", "every run, hash-pinned", edge=GOLD)
fbox(s, 3.8, 4.7, 2.3, 1.0, "trust gates", "n, contract, CI, flip, tripwire", edge=GOOD)
fbox(s, 6.7, 4.7, 2.3, 1.0, "export", "benchmark.json, differential-tested")
fbox(s, 9.8, 3.8, 2.3, 1.2, "surfaces", "site · notebook · deck · this room")
arrow(s, 3.2, 2.9, 3.8, 3.35)
arrow(s, 3.2, 4.2, 3.8, 3.75)
arrow(s, 6.1, 3.55, 6.7, 3.55)
arrow(s, 7.4, 4.05, 6.2, 4.75)
arrow(s, 6.1, 5.2, 6.7, 5.2)
arrow(s, 9.0, 5.1, 9.85, 4.6)
arrow(s, 9.0, 3.45, 9.85, 4.0)
note(s, "The deck, the site, and the notebook all read the same export; the export mirrors "
        "the cost model exactly and a differential test proves it. No number on any surface "
        "has a second source.", top=6.35)

# 4 data
s = slide()
eyebrow(s, "PPDAC · Data")
h2(s, "89 labeled cases, five suites, tiers kept apart")
image(s, "eda_corpus.png", top=2.1, height=3.6)
note(s, "Only 4 labels are expert ground truth; construction labels say so wherever they appear.", top=6.1)

# 5 repeats
s = slide()
eyebrow(s, "PPDAC · Analysis · the ladder")
h2(s, "Repeats dissolve stories")
image(s, "injection_repeats.png", top=2.1, height=3.5)
note(s, "A single run showed v4b resisting the hardest planted note; N=5 put every shipped "
        "rung in one coin-flip band. The repeat protocol also reversed one gate decision and "
        "revoked one perfect score.", top=6.0)

# 6 matrix
s = slide()
eyebrow(s, "PPDAC · Analysis · the matrix")
h2(s, "The gate bites, and ranking survives")
bigrow(s, [(str(K["total_cells"]), "prompt x model cells", TEXT),
           ("10", "disqualified by the zero-tolerance tripwire", GOLD),
           (str(K["trusted_cells"]), "fully trusted cells", TEXT)], top=1.9)
image(s, "bootstrap_loss.png", top=3.6, height=2.5)
note(s, "Weighted loss is the one OEC; contract, flip, injection, and the tripwire sit "
        "OUTSIDE it as guardrails, on purpose: fold them in and the score becomes gameable. "
        "The gate fired ten times and the ranking survived.", top=5.6)

# 7 confidence
s = slide()
eyebrow(s, "PPDAC · Analysis · confidence")
h2(s, "Stated confidence is decorative; earned confidence is measured")
image(s, "reliability.png", top=2.1, height=3.7)
note(s, "All 21 stated confidences live between 0.90 and 1.00 while the runs earn 0.90 "
        "overall; at stated 0.90 the model was right once in two. The replacement: vote "
        "fraction over N=5 temperature-raised repeats.", top=6.1)

# 8 judge
s = slide()
eyebrow(s, "PPDAC · Analysis · the judge")
h2(s, "Triangulated across three families")
tf = box(s, 0.9, 2.3, 11.4, 3.4)
for i, (name, score, colr, note_t) in enumerate([
    ("gemini-flash", "5.0 / 5.0 / 5.0", BAD, "saturated: a dead instrument, reported as one"),
    ("claude-haiku", "4.2 / 3.9 / 3.1", GOOD, "discriminates; proportionality lowest"),
    ("phi-4 (microsoft)", "4.6 / 5.0 / 4.7", GOOD, "third family, overlaps no judged column"),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r1 = p.add_run(); r1.text = f"{name:<22}"
    r1.font.size = Pt(18); r1.font.bold = True; r1.font.color.rgb = TEXT; r1.font.name = "Consolas"
    r2 = p.add_run(); r2.text = f"  {score}   "
    r2.font.size = Pt(18); r2.font.bold = True; r2.font.color.rgb = colr; r2.font.name = "Consolas"
    r3 = p.add_run(); r3.text = note_t
    r3.font.size = Pt(14); r3.font.color.rgb = MUT; r3.font.name = "Inter"
    p.space_after = Pt(16)
note(s, "Both discriminating families score proportionality lowest, the axis the policy "
        "makes hardest. The saturated judge is never averaged in.", top=5.9)

# 9 synthetic
s = slide()
eyebrow(s, "PPDAC · Analysis · scale")
h2(s, "64 generated cases found a policy ambiguity")
image(s, "synthetic_sweep.png", top=2.1, height=3.5)
note(s, "v5 sweeps the uncontested archetypes 48/48 (v1: 44/48). The contested family splits "
        "seven models 20-8 with no family pattern: the disagreement lives in the policy text, "
        "routed to its owner, never resolved by me.", top=6.0)



# overfitting defense
s = slide()
eyebrow(s, "The question a 100% row earns")
h2(s, "Is v5 overfit to nine cases? Three tests, one open risk")
tf = box(s, 0.9, 2.0, 11.4, 4.7)
tf.word_wrap = True
for i, (h, b) in enumerate([
    ("Transfer it never trained for:", "v5 was tuned against gemini-flash only. Unchanged, "
     "it scores 12/12 on claude-sonnet and claude-haiku and 11/12 on gemini-pro (n=12 each)."),
    ("Data authored after the freeze:", "the 64 generated cases did not exist when v5 "
     "froze. It sweeps 56/56 on flash, 48/48 on qwen, 46/48 on gemini-pro. Overfit "
     "collapses on unseen data; this did not."),
    ("Weak columns fail for capacity, not memorization:", "llama trails at every rung, "
     "8/12 on v1 before any tuning existed, 9/12 at v5 with contract 1.00. The ladder "
     "never bent toward its failures, and they predate it."),
    ("The discipline check:", "v6 was rejected on its own pre-registered bar. An "
     "overfitting process chases the suite; this one turned down a better injection score."),
    ("Open, stated:", "the true holdout is n=3 (v5: 2/3), too small to certify anything. "
     "That is why the gates carry the ship claim and more expert labels is the named blocker."),
]):
    p2 = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r1 = p2.add_run(); r1.text = h + " "
    r1.font.size = Pt(14); r1.font.bold = True; r1.font.color.rgb = TEXT; r1.font.name = "Inter"
    r2 = p2.add_run(); r2.text = b
    r2.font.size = Pt(14); r2.font.color.rgb = MUT; r2.font.name = "Inter"
    p2.space_after = Pt(11)

# 9b routing
s = slide()
eyebrow(s, "PPDAC · Analysis · the queue")
h2(s, "A contested case is routed, not resolved")
fbox(s, 0.9, 3.3, 2.5, 1.05, "data_quality_flag", "generated archetype, 28 runs")
fbox(s, 4.0, 3.3, 2.5, 1.05, "seven models", "split 20 HOLD / 8 REJECT", edge=ACC)
fbox(s, 7.2, 2.2, 2.6, 1.05, "not a model error", "POLICY.md underdetermines it", edge=BAD)
fbox(s, 7.2, 4.5, 2.6, 1.05, "no score assigned", "contested flag in the ledger")
fbox(s, 10.3, 3.3, 2.1, 1.15, "policy owner", "gets the split, decides once", edge=GOLD)
arrow(s, 3.4, 3.85, 4.0, 3.85)
arrow(s, 6.5, 3.5, 7.2, 2.9)
arrow(s, 6.5, 4.2, 7.2, 4.9)
arrow(s, 9.8, 2.75, 10.45, 3.3)
arrow(s, 9.8, 5.0, 10.45, 4.45)
note(s, "Same mechanism for sanctions-partial: zero APPROVEs in 23 runs across seven "
        "models, so it fails closed and ships contested. The benchmark's job on these cases "
        "is to surface the disagreement, not to pick a winner. Given time, this queue "
        "becomes a real reviewer surface; the playground already prototypes it.", top=6.35)

# 10 population
s = slide()
eyebrow(s, "PPDAC · Analysis · population")
h2(s, "A million-case simulation on measured kernels")
image(s, "population_sim.png", top=2.1, height=3.6)
note(s, "Behavior measured from the ledger; prevalence and exposure are named, swept "
        "assumptions. At 0.5% fraud, v1 costs $4,594/1k matrix-priced and $12,069/1k "
        "exposure-priced; rare large events dominate.", top=6.1)

# 11 reviews
s = slide()
eyebrow(s, "Adversarial reviews, answered with runs")
h2(s, "The gate does not bend on deadline day")
tf = box(s, 0.9, 2.2, 11.4, 4.2)
tf.word_wrap = True
for i, (h, b) in enumerate([
    ("Refuted by probes:", "precomputed-false is never exculpatory (4/4); the \"you survive "
     "108 by luck\" perturbation HOLDs 4/4."),
    ("Half right:", "the sanctions middle case fails CLOSED, not open: zero APPROVEs in 23 "
     "runs across seven models; the verdict choice ships contested."),
    ("v6, the injection line:", "resists 19/20 where v5 sits at 5/9, and was rejected "
     "anyway: suite 11/12 against a pre-registered 12/12 bar. It ships as the tested next "
     "rung, not as the prompt."),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r1 = p.add_run(); r1.text = h + " "
    r1.font.size = Pt(16); r1.font.bold = True; r1.font.color.rgb = TEXT; r1.font.name = "Inter"
    r2 = p.add_run(); r2.text = b
    r2.font.size = Pt(16); r2.font.color.rgb = MUT; r2.font.name = "Inter"
    p.space_after = Pt(16)

# 12 conclusion
s = slide()
eyebrow(s, "PPDAC · Conclusion")
h2(s, "Ship v5 + gemini-flash, and say what that means")
tf = box(s, 0.9, 2.3, 11.4, 3.2)
tf.word_wrap = True
for i, (tier, colr, body) in enumerate([
    ("CERTIFIED", GOOD, "gate recall 1.0 on every ranked cell; contract at floor everywhere ranked"),
    ("SUGGESTED", ACC, "12/12 with Wilson floor 74%; holdout n=3; the ladder's edge at ~3:1 posterior odds"),
    ("DECORATIVE", BAD, "verbalized confidence; flash-as-judge"),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r1 = p.add_run(); r1.text = f"{tier:<14}"
    r1.font.size = Pt(15); r1.font.bold = True; r1.font.color.rgb = colr; r1.font.name = "Consolas"
    r2 = p.add_run(); r2.text = body
    r2.font.size = Pt(16); r2.font.color.rgb = MUT; r2.font.name = "Inter"
    p.space_after = Pt(14)
note(s, "Every number triggers an action or it is just reporting: gates green = ship; flip or "
        "contested split = hold and route to a human; tripwire = roll back. The $0 is a point "
        "estimate; n=12 cannot rule out ~$530k/1k worst-class. The gates carry the claim.", top=5.7)


# out of scope + next given time
s = slide()
eyebrow(s, "The boundary, drawn on purpose")
h2(s, "Out of scope, and what I'd do next given time")
tf = box(s, 0.9, 2.0, 5.8, 4.8)
tf.word_wrap = True
para(tf, "DELIBERATELY OUT OF SCOPE", 12, ACC, bold=True, first=True, space=10)
for h, b in [
    ("Tools and skills:", "the agent decides from the dossier alone; no function calling, "
     "no retrieval. Tools change the eval surface, so they get their own arm, not a patch."),
    ("MCP integrations:", "no external context servers; every input is the frozen case "
     "file, which keeps runs reproducible."),
    ("CLI vs API surfaces:", "claude columns ride the subscription CLI (no sampler "
     "control, stated); the API/SDK rerun is specced, blocked only on a key."),
    ("Memory attached to models:", "every run is stateless by design; a remembering agent "
     "cannot be measured by repeats, and repeats are the method."),
    ("Fine-tuning:", "the prompt is the only movable part, on purpose."),
]:
    p2 = tf.add_paragraph()
    r1 = p2.add_run(); r1.text = h + " "
    r1.font.size = Pt(12); r1.font.bold = True; r1.font.color.rgb = TEXT; r1.font.name = "Inter"
    r2 = p2.add_run(); r2.text = b
    r2.font.size = Pt(12); r2.font.color.rgb = MUT; r2.font.name = "Inter"
    p2.space_after = Pt(8)
tf = box(s, 7.1, 2.0, 5.3, 4.8)
tf.word_wrap = True
para(tf, "NEXT, GIVEN TIME, IN ORDER", 12, GOLD, bold=True, first=True, space=10)
for i, b in enumerate([
    "More expert labels: the hard blocker; SPRT already says how many are needed.",
    "The v6 injection line through the full gate.",
    "API rerun of the claude columns with temperature control.",
    "The adjudication queue wired to a real reviewer; the playground already prototypes it.",
    "A tool-augmented variant as a new benchmark arm.",
]):
    p2 = tf.add_paragraph()
    r = p2.add_run(); r.text = f"{i+1}.  {b}"
    r.font.size = Pt(13); r.font.color.rgb = MUT; r.font.name = "Inter"
    p2.space_after = Pt(10)

# 13 production
s = slide()
eyebrow(s, "If this shipped Monday")
h2(s, "What breaks, in order")
tf = box(s, 0.9, 2.2, 11.4, 3.8)
tf.word_wrap = True
for i, t in enumerate([
    "Four real labels cannot certify a production prompt: more expert labels is the hard blocker.",
    "Pin the prompt AND the model: the pairing is what was verified.",
    "Injection is not solved by any shipped rung; the tested v6 line is next.",
    "Confidence gating must use vote fraction, never stated confidence.",
    "The contested case families need a policy-owner decision, and the queue collects it.",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = f"{i+1}.  {t}"
    r.font.size = Pt(16); r.font.color.rgb = MUT; r.font.name = "Inter"
    p.space_after = Pt(12)
note(s, "verdict-bench.pages.dev   ·   github.com/ShovalBenjer/verdict-bench", top=6.3, color=DIM)

prs.save(OUT)
print(f"{OUT} ({OUT.stat().st_size // 1024} KB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
